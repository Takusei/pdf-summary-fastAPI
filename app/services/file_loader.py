import os
from typing import List

import pandas as pd
from langchain_community.document_loaders import (
    Docx2txtLoader,
    PyMuPDFLoader,
    # UnstructuredExcelLoader, # This need unstructured package, but python-magic is not available in all environments
    # UnstructuredPowerPointLoader, # This need system level dependencies
)
from langchain_core.documents import Document
from pptx import Presentation


def load_pdf(file_path: str) -> list[Document]:
    loader = PyMuPDFLoader(file_path)
    docs = loader.load()
    if not docs:
        raise ValueError("PDF is empty or unreadable")
    return docs


def load_pptx(file_path: str) -> list[Document]:
    # loader = UnstructuredPowerPointLoader(file_path)
    # docs = loader.load()
    # if not docs:
    #     raise ValueError("PPTX is empty or unreadable")
    # return docs
    prs = Presentation(file_path)
    docs: List[Document] = []

    for i, slide in enumerate(prs.slides, start=1):
        parts = []

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = shape.text.strip()
                if t:
                    parts.append(t)

        if not parts:
            continue

        text = "\n".join(parts)

        title = None
        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text.strip()

        docs.append(
            Document(
                page_content=text,
                metadata={
                    "source": file_path,
                    "slide": i,
                    "title": title,
                },
            )
        )
    return docs


def load_docx(file_path: str) -> list[Document]:
    loader = Docx2txtLoader(file_path)
    docs = loader.load()
    if not docs:
        raise ValueError("DOCX is empty or unreadable")
    return docs


def load_excel(file_path: str) -> list[Document]:
    try:
        sheets = pd.read_excel(file_path, sheet_name=None)
    except Exception as exc:
        raise ValueError("Excel file is empty or unreadable") from exc

    docs: List[Document] = []
    for sheet_name, df in sheets.items():
        if df is None or df.empty:
            continue
        text = f"Sheet: {sheet_name}\n" + df.to_csv(index=False)
        docs.append(
            Document(
                page_content=text,
                metadata={
                    "source": file_path,
                    "sheet": sheet_name,
                },
            )
        )

    if not docs:
        raise ValueError("Excel file is empty or unreadable")

    return docs


# Create a dispatch table for file loaders
FILE_LOADERS = {
    ".pdf": load_pdf,
    ".pptx": load_pptx,
    # ".ppt": load_pptx,  # .ppt is not supported by python-pptx
    ".docx": load_docx,
    ".xls": load_excel,
    ".xlsx": load_excel,
}


def load_file(file_path: str) -> list[Document]:
    """
    Loads a file using the appropriate loader based on its extension.
    """
    _, extension = os.path.splitext(file_path)
    extension = extension.lower()

    loader_func = FILE_LOADERS.get(extension)
    if not loader_func:
        raise ValueError(f"Unsupported file type: {extension}")

    return loader_func(file_path)
